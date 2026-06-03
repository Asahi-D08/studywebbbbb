"""Build a StudyWeb presentation (.pptx) — Keynote opens it natively.

Run:
    .venv/bin/python scripts/build_studyweb_keynote.py
Output:
    StudyWeb.pptx  (double-click to open in Keynote)
"""

from __future__ import annotations

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.oxml.ns import qn
from pptx.util import Emu, Inches, Pt
from lxml import etree


# ─── Theme ────────────────────────────────────────────────────────────────────

BG_DARK = RGBColor(0x0B, 0x0F, 0x1E)
BG_PANEL = RGBColor(0x16, 0x1B, 0x2E)
ACCENT_INDIGO = RGBColor(0x63, 0x66, 0xF1)
ACCENT_VIOLET = RGBColor(0x8B, 0x5C, 0xF6)
ACCENT_FUCHSIA = RGBColor(0xD9, 0x46, 0xEF)
TEXT_PRIMARY = RGBColor(0xF8, 0xFA, 0xFC)
TEXT_MUTED = RGBColor(0x94, 0xA3, 0xB8)
TEXT_DIM = RGBColor(0x64, 0x74, 0x8B)
LINE = RGBColor(0x1E, 0x29, 0x3B)

FONT = "Helvetica Neue"
FONT_MONO = "Menlo"


# ─── Helpers ──────────────────────────────────────────────────────────────────

def add_bg(slide, prs):
    """Solid dark background + a subtle gradient bar at the top."""
    w, h = prs.slide_width, prs.slide_height
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, w, h)
    bg.line.fill.background()
    bg.fill.solid()
    bg.fill.fore_color.rgb = BG_DARK
    bg.shadow.inherit = False

    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, w, Inches(0.08))
    bar.line.fill.background()
    _set_gradient_fill(bar, ACCENT_INDIGO, ACCENT_VIOLET, ACCENT_FUCHSIA)
    bar.shadow.inherit = False


def _set_gradient_fill(shape, c1: RGBColor, c2: RGBColor, c3: RGBColor) -> None:
    """Apply a 3-stop horizontal gradient by writing the spPr XML directly."""
    spPr = shape.fill._xPr  # this is already <a:spPr> for shapes
    for child in list(spPr):
        if child.tag in (qn("a:noFill"), qn("a:solidFill"), qn("a:gradFill"),
                         qn("a:blipFill"), qn("a:pattFill")):
            spPr.remove(child)

    grad = etree.SubElement(spPr, qn("a:gradFill"), {"flip": "none", "rotWithShape": "1"})
    lst = etree.SubElement(grad, qn("a:gsLst"))
    for pos, color in [(0, c1), (50000, c2), (100000, c3)]:
        gs = etree.SubElement(lst, qn("a:gs"), {"pos": str(pos)})
        etree.SubElement(
            gs,
            qn("a:srgbClr"),
            {"val": f"{color[0]:02X}{color[1]:02X}{color[2]:02X}"},
        )
    etree.SubElement(grad, qn("a:lin"), {"ang": "0", "scaled": "1"})


def add_text(
    slide,
    text: str,
    left,
    top,
    width,
    height,
    *,
    size: int = 18,
    bold: bool = False,
    color: RGBColor = TEXT_PRIMARY,
    font: str = FONT,
    align: PP_ALIGN = PP_ALIGN.LEFT,
    anchor: MSO_ANCHOR = MSO_ANCHOR.TOP,
    spacing: float = 1.15,
):
    tb = slide.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = 0
    tf.margin_top = tf.margin_bottom = 0
    tf.vertical_anchor = anchor

    lines = text.split("\n") if isinstance(text, str) else list(text)
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        p.line_spacing = spacing
        r = p.add_run()
        r.text = line
        r.font.name = font
        r.font.size = Pt(size)
        r.font.bold = bold
        r.font.color.rgb = color
    return tb


def add_eyebrow(slide, text: str, left, top):
    """Small uppercase tracked label above titles."""
    add_text(
        slide,
        text.upper(),
        left,
        top,
        Inches(6),
        Inches(0.3),
        size=11,
        bold=True,
        color=TEXT_MUTED,
    )


def add_page_number(slide, n: int, total: int, prs):
    add_text(
        slide,
        f"{n:02d} / {total:02d}",
        prs.slide_width - Inches(1.2),
        prs.slide_height - Inches(0.55),
        Inches(1),
        Inches(0.3),
        size=10,
        color=TEXT_DIM,
        align=PP_ALIGN.RIGHT,
        font=FONT_MONO,
    )


def add_footer(slide, prs):
    add_text(
        slide,
        "StudyWeb · Notes Hub",
        Inches(0.6),
        prs.slide_height - Inches(0.55),
        Inches(4),
        Inches(0.3),
        size=10,
        color=TEXT_DIM,
    )


def add_bullets(slide, items, left, top, width, height, *, size: int = 18,
                gap: float = 1.55, bullet_color: RGBColor = ACCENT_VIOLET):
    """A list with colored squares as bullets."""
    tb = slide.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = 0
    tf.margin_top = tf.margin_bottom = 0

    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = PP_ALIGN.LEFT
        p.line_spacing = gap
        p.space_after = Pt(6)

        bullet = p.add_run()
        bullet.text = "■  "
        bullet.font.name = FONT
        bullet.font.size = Pt(size)
        bullet.font.color.rgb = bullet_color

        if isinstance(item, tuple):
            head, tail = item
            r1 = p.add_run()
            r1.text = head
            r1.font.name = FONT
            r1.font.size = Pt(size)
            r1.font.bold = True
            r1.font.color.rgb = TEXT_PRIMARY

            r2 = p.add_run()
            r2.text = "  " + tail
            r2.font.name = FONT
            r2.font.size = Pt(size)
            r2.font.color.rgb = TEXT_MUTED
        else:
            r = p.add_run()
            r.text = item
            r.font.name = FONT
            r.font.size = Pt(size)
            r.font.color.rgb = TEXT_PRIMARY


def add_panel(slide, left, top, width, height, *, fill=BG_PANEL,
              line=LINE, radius: bool = True):
    shape_type = MSO_SHAPE.ROUNDED_RECTANGLE if radius else MSO_SHAPE.RECTANGLE
    s = slide.shapes.add_shape(shape_type, left, top, width, height)
    s.fill.solid()
    s.fill.fore_color.rgb = fill
    s.line.color.rgb = line
    s.line.width = Pt(0.75)
    s.shadow.inherit = False
    if radius and s.adjustments:
        s.adjustments[0] = 0.04
    return s


# ─── Slides ───────────────────────────────────────────────────────────────────

TOTAL = 6


def slide_cover(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide, prs)

    add_eyebrow(slide, "Project introduction · 2026", Inches(0.6), Inches(0.6))

    add_text(
        slide, "StudyWeb",
        Inches(0.6), Inches(1.4), Inches(12), Inches(1.6),
        size=84, bold=True, color=TEXT_PRIMARY,
    )

    # Gradient subtitle line
    gradient_strip = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(0.6), Inches(3.1), Inches(2.2), Inches(0.08)
    )
    gradient_strip.line.fill.background()
    _set_gradient_fill(gradient_strip, ACCENT_INDIGO, ACCENT_VIOLET, ACCENT_FUCHSIA)
    gradient_strip.shadow.inherit = False

    add_text(
        slide,
        "A notes hub built for IGCSE & IB students.",
        Inches(0.6), Inches(3.45), Inches(12), Inches(0.7),
        size=28, color=TEXT_PRIMARY,
    )
    add_text(
        slide,
        "Tech stack · Implementation · What sets it apart",
        Inches(0.6), Inches(4.05), Inches(12), Inches(0.6),
        size=20, color=TEXT_MUTED,
    )

    add_text(
        slide,
        "Next.js 16   ·   React 19   ·   TypeScript 5   ·   Tailwind v4",
        Inches(0.6), Inches(6.5), Inches(12), Inches(0.4),
        size=13, color=TEXT_DIM, font=FONT_MONO,
    )
    add_page_number(slide, 1, TOTAL, prs)


def slide_overview(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide, prs)
    add_eyebrow(slide, "01 · Overview", Inches(0.6), Inches(0.6))
    add_text(
        slide, "What is StudyWeb?",
        Inches(0.6), Inches(0.95), Inches(12), Inches(0.9),
        size=40, bold=True,
    )

    add_text(
        slide,
        "A lightweight web app where students browse, upload and discuss\n"
        "study notes organised by programme, section and subject.",
        Inches(0.6), Inches(1.95), Inches(12), Inches(1.2),
        size=20, color=TEXT_MUTED, spacing=1.3,
    )

    add_bullets(
        slide,
        [
            ("Programme-aware",   "Built around real IGCSE and IB course structures."),
            ("Three core modules", "Public Sharing area · Private notes · Per-section Forum."),
            ("Open to visitors",   "Read shared notes without signing up; sign in to contribute."),
            ("Clean URL design",   "/[program]/[section]/[subject] — predictable and shareable."),
        ],
        Inches(0.6), Inches(3.6), Inches(12), Inches(3.5),
        size=18,
    )

    add_footer(slide, prs)
    add_page_number(slide, 2, TOTAL, prs)


def slide_stack(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide, prs)
    add_eyebrow(slide, "02 · Technology", Inches(0.6), Inches(0.6))
    add_text(
        slide, "Tech stack",
        Inches(0.6), Inches(0.95), Inches(12), Inches(0.9),
        size=40, bold=True,
    )
    add_text(
        slide,
        "A modern full-stack TypeScript app — under a dozen runtime dependencies.",
        Inches(0.6), Inches(1.95), Inches(12), Inches(0.6),
        size=18, color=TEXT_MUTED,
    )

    # 2-column card grid
    rows = [
        ("Framework",  "Next.js 16.2.1  ·  App Router + RSC"),
        ("UI",         "React 19  ·  Tailwind CSS v4"),
        ("Language",   "TypeScript 5  ·  ESLint 9"),
        ("Auth",       "jose JWT  ·  bcryptjs hashing"),
        ("Theme",      "next-themes  ·  light + dark"),
        ("Storage",    "Local JSON + flat file dir  ·  no database"),
    ]

    col_w = Inches(6.0)
    card_h = Inches(0.9)
    gap_y = Inches(0.2)
    gap_x = Inches(0.3)
    start_x = Inches(0.6)
    start_y = Inches(2.85)

    for i, (label, value) in enumerate(rows):
        col = i % 2
        row = i // 2
        x = start_x + (col_w + gap_x) * col
        y = start_y + (card_h + gap_y) * row
        add_panel(slide, x, y, col_w, card_h)
        add_text(slide, label.upper(),
                 x + Inches(0.3), y + Inches(0.13),
                 col_w - Inches(0.6), Inches(0.3),
                 size=11, bold=True, color=ACCENT_VIOLET)
        add_text(slide, value,
                 x + Inches(0.3), y + Inches(0.42),
                 col_w - Inches(0.6), Inches(0.4),
                 size=16, color=TEXT_PRIMARY)

    add_footer(slide, prs)
    add_page_number(slide, 3, TOTAL, prs)


def slide_implementation(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide, prs)
    add_eyebrow(slide, "03 · Implementation", Inches(0.6), Inches(0.6))
    add_text(
        slide, "How it's built",
        Inches(0.6), Inches(0.95), Inches(12), Inches(0.9),
        size=40, bold=True,
    )

    # Two columns: Architecture | Data & Security
    col_w = Inches(6.0)
    col_h = Inches(4.6)
    gap = Inches(0.3)

    # Left card
    add_panel(slide, Inches(0.6), Inches(2.1), col_w, col_h)
    add_text(slide, "HOW THE PAGES WORK",
             Inches(0.95), Inches(2.3), col_w - Inches(0.7), Inches(0.4),
             size=12, bold=True, color=ACCENT_INDIGO)
    add_text(slide, "Cook on the server, ship a ready meal.",
             Inches(0.95), Inches(2.6), col_w - Inches(0.7), Inches(0.5),
             size=20, bold=True)
    add_bullets(
        slide,
        [
            "Pages are rendered on the server — you get HTML, not a loading spinner.",
            "Each page grabs its own data; no extra trips back to the backend.",
            "Folder names map straight to URLs — easy to add a new subject.",
            "A single gatekeeper checks every request before it gets in.",
            "Core logic lives in one place, so storage can be swapped later.",
        ],
        Inches(0.95), Inches(3.3), col_w - Inches(0.7), Inches(3.2),
        size=13, gap=1.3,
    )

    # Right card
    right_x = Inches(0.6) + col_w + gap
    add_panel(slide, right_x, Inches(2.1), col_w, col_h)
    add_text(slide, "HOW DATA STAYS SAFE",
             right_x + Inches(0.35), Inches(2.3), col_w - Inches(0.7), Inches(0.4),
             size=12, bold=True, color=ACCENT_FUCHSIA)
    add_text(slide, "No database, no half-saved files, no leaks.",
             right_x + Inches(0.35), Inches(2.6), col_w - Inches(0.7), Inches(0.5),
             size=20, bold=True)
    add_bullets(
        slide,
        [
            "Notes live in plain JSON files — back up by zipping one folder.",
            "Saves write to a temp file first, so a crash never breaks data.",
            "Saves run one at a time, so users can't overwrite each other.",
            "Passwords are scrambled; login tokens are signed and tamper-proof.",
            "Uploaded filenames are cleaned to block path-tricks.",
        ],
        right_x + Inches(0.35), Inches(3.3), col_w - Inches(0.7), Inches(3.2),
        size=13, gap=1.3, bullet_color=ACCENT_FUCHSIA,
    )

    add_footer(slide, prs)
    add_page_number(slide, 4, TOTAL, prs)


def slide_comparison(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide, prs)
    add_eyebrow(slide, "04 · Comparison", Inches(0.6), Inches(0.6))
    add_text(
        slide, "How it compares",
        Inches(0.6), Inches(0.95), Inches(12), Inches(0.9),
        size=40, bold=True,
    )
    add_text(
        slide,
        "Against the tools students currently mix together.",
        Inches(0.6), Inches(1.95), Inches(12), Inches(0.5),
        size=16, color=TEXT_MUTED,
    )

    headers = ["Capability", "StudyWeb", "Google Classroom", "Wenku / CSDN"]
    rows = [
        ("Course-structured navigation", "Yes", "Class-based", "No"),
        ("Public + private per note",    "Yes", "No", "No"),
        ("One-click copy as snapshot",   "Yes", "No", "No"),
        ("Subject forum + note citing",  "Yes", "Yes", "Comments"),
        ("Self-hosted, data you own",    "Yes", "No", "No"),
        ("Zero-DB deployment",           "Yes", "—", "—"),
    ]

    table_left = Inches(0.6)
    table_top = Inches(2.65)
    table_w = Inches(12.4)
    col_widths = [Inches(4.0), Inches(2.5), Inches(3.0), Inches(2.9)]
    row_h = Inches(0.55)

    # Header row
    x = table_left
    for i, h in enumerate(headers):
        cell_fill = ACCENT_INDIGO if i == 1 else BG_PANEL
        cell = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, table_top, col_widths[i], row_h)
        cell.fill.solid(); cell.fill.fore_color.rgb = cell_fill
        cell.line.color.rgb = LINE; cell.line.width = Pt(0.5)
        cell.shadow.inherit = False
        add_text(
            slide, h,
            x + Inches(0.2), table_top + Inches(0.12),
            col_widths[i] - Inches(0.4), row_h - Inches(0.2),
            size=12, bold=True,
            color=TEXT_PRIMARY,
        )
        x += col_widths[i]

    # Body rows
    for r_idx, row in enumerate(rows):
        y = table_top + row_h + row_h * r_idx
        x = table_left
        zebra = BG_PANEL if r_idx % 2 == 0 else RGBColor(0x10, 0x16, 0x29)
        for c_idx, val in enumerate(row):
            cell = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, col_widths[c_idx], row_h)
            cell.fill.solid()
            cell.fill.fore_color.rgb = zebra
            cell.line.color.rgb = LINE
            cell.line.width = Pt(0.5)
            cell.shadow.inherit = False

            if c_idx == 0:
                color = TEXT_PRIMARY
                bold = True
            elif c_idx == 1 and val == "Yes":
                color = ACCENT_VIOLET
                bold = True
            elif val == "Yes":
                color = TEXT_PRIMARY
                bold = False
            elif val in ("No", "—"):
                color = TEXT_DIM
                bold = False
            else:
                color = TEXT_MUTED
                bold = False
            add_text(
                slide, val,
                x + Inches(0.2), y + Inches(0.13),
                col_widths[c_idx] - Inches(0.4), row_h - Inches(0.2),
                size=12, bold=bold, color=color,
            )
            x += col_widths[c_idx]

    add_footer(slide, prs)
    add_page_number(slide, 5, TOTAL, prs)


def slide_ux(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide, prs)
    add_eyebrow(slide, "05 · Why it stands out", Inches(0.6), Inches(0.6))
    add_text(
        slide, "Designed around the student experience",
        Inches(0.6), Inches(0.95), Inches(12), Inches(0.9),
        size=36, bold=True,
    )

    items = [
        ("Browse without friction",
         "No sign-up wall — visitors read every shared note instantly."),
        ("Course is the navigation",
         "Three clicks: IGCSE/IB → section → subject. No search needed."),
        ("Public and private, one tap apart",
         "Toggle any note between Sharing area and Private at any time."),
        ("Save the good ones forever",
         "Copy a shared note into your private library as an independent snapshot."),
        ("Discuss right next to the notes",
         "Each section has a forum — quote notes or attach files in one composer."),
        ("Polished little details",
         "Drag-and-drop upload, instant validation, dark mode, mobile-ready."),
        ("Quiet by design",
         "No ads, no notifications, no tracking — just your notes."),
    ]

    col_w = Inches(6.0)
    card_h = Inches(1.25)
    gap_x = Inches(0.3)
    gap_y = Inches(0.2)
    start_x = Inches(0.6)
    start_y = Inches(2.05)

    for i, (head, body) in enumerate(items):
        col = i % 2
        row = i // 2
        x = start_x + (col_w + gap_x) * col
        y = start_y + (card_h + gap_y) * row
        add_panel(slide, x, y, col_w, card_h)
        # tiny accent dot
        dot = slide.shapes.add_shape(MSO_SHAPE.OVAL,
                                     x + Inches(0.3), y + Inches(0.28),
                                     Inches(0.18), Inches(0.18))
        dot.line.fill.background()
        dot.fill.solid()
        dot.fill.fore_color.rgb = ACCENT_FUCHSIA if i % 2 else ACCENT_INDIGO
        dot.shadow.inherit = False

        add_text(slide, head,
                 x + Inches(0.65), y + Inches(0.18),
                 col_w - Inches(0.85), Inches(0.4),
                 size=15, bold=True, color=TEXT_PRIMARY)
        add_text(slide, body,
                 x + Inches(0.65), y + Inches(0.55),
                 col_w - Inches(0.85), Inches(0.7),
                 size=12, color=TEXT_MUTED, spacing=1.25)

    add_footer(slide, prs)
    add_page_number(slide, 6, TOTAL, prs)


# ─── Build ────────────────────────────────────────────────────────────────────

def build(output_path: str) -> None:
    prs = Presentation()
    # 16:9 — 13.333 × 7.5 inches
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    slide_cover(prs)
    slide_overview(prs)
    slide_stack(prs)
    slide_implementation(prs)
    slide_comparison(prs)
    slide_ux(prs)

    prs.save(output_path)


if __name__ == "__main__":
    import os
    out = os.path.join(os.path.dirname(os.path.dirname(os.path.abspath(__file__))),
                       "StudyWeb.pptx")
    build(out)
    print(f"Wrote {out}")
